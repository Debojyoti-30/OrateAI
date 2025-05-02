import os
from pptx import Presentation
import fitz  # PyMuPDF
from groq import Groq

# Function to extract text from PPTX
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text_runs = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_text = f"Slide {slide_number}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text.strip() + "\n"
        text_runs.append(slide_text.strip())
    return "\n\n".join(text_runs)

# Function to extract text from PDF
def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    pdf_text = []
    for page in doc:
        pdf_text.append(page.get_text())
    return "\n\n".join(pdf_text)

# Function to summarize using Groq API
def summarize_with_groq(content):
    client = Groq(
        api_key=os.environ.get("GROQ_API_KEY")  # Set securely as env var
    )
    chat_completion = client.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": f"Summarize the content of this document:\n{content}",
            }
        ],
        model="llama-3.3-70b-versatile",
        stream=False,
    )
    return chat_completion.choices[0].message.content

# ------------------------------
# ✅ Usage Example (Jupyter/Colab)
# ------------------------------
file_path = "BPlan-1-17.pdf"  # or "yourfile.pdf"
file_type = "pdf"        # or "pdf"

# Extract text
if file_type == "pptx":
    extracted_text = extract_text_from_pptx(file_path)
elif file_type == "pdf":
    extracted_text = extract_text_from_pdf(file_path)
else:
    raise ValueError("Unsupported file type")

# Summarize
summary = summarize_with_groq(extracted_text)
print(summary)
